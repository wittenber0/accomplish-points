#!/usr/bin/env python3
"""Extract all text content and images from the PowerPoint file."""

import os
import json
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

PPTX_PATH = "WEBSITE DRAFT MATERIALS  ACCOMPLISH POINTS .pptx"
IMAGES_DIR = "extracted_images"

os.makedirs(IMAGES_DIR, exist_ok=True)

prs = Presentation(PPTX_PATH)

slides_data = []
image_count = 0

for slide_idx, slide in enumerate(prs.slides, 1):
    slide_info = {
        "slide_number": slide_idx,
        "layout_name": slide.slide_layout.name if slide.slide_layout else "Unknown",
        "shapes": [],
        "notes": ""
    }

    # Extract notes
    if slide.has_notes_slide:
        notes_slide = slide.notes_slide
        notes_text = notes_slide.notes_text_frame.text if notes_slide.notes_text_frame else ""
        slide_info["notes"] = notes_text

    for shape in slide.shapes:
        shape_info = {
            "name": shape.name,
            "shape_type": str(shape.shape_type),
            "position": {
                "left": shape.left,
                "top": shape.top,
                "width": shape.width,
                "height": shape.height
            }
        }

        # Text content
        if shape.has_text_frame:
            paragraphs = []
            for para in shape.text_frame.paragraphs:
                runs_text = []
                for run in para.runs:
                    try:
                        color = str(run.font.color.rgb) if run.font.color and run.font.color.type is not None else None
                    except (AttributeError, TypeError):
                        color = None
                    run_info = {
                        "text": run.text,
                        "bold": run.font.bold,
                        "italic": run.font.italic,
                        "size": str(run.font.size) if run.font.size else None,
                        "color": color
                    }
                    runs_text.append(run_info)
                paragraphs.append({
                    "text": para.text,
                    "level": para.level,
                    "runs": runs_text
                })
            shape_info["text_content"] = paragraphs

        # Table content
        if shape.has_table:
            table_data = []
            for row in shape.table.rows:
                row_data = []
                for cell in row.cells:
                    row_data.append(cell.text)
                table_data.append(row_data)
            shape_info["table"] = table_data

        # Image content
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image = shape.image
            image_count += 1
            try:
                ext = image.content_type.split("/")[-1]
                if ext == "jpeg":
                    ext = "jpg"
            except (ValueError, AttributeError):
                # Fallback for unsupported formats
                blob = image.blob
                if blob[:4] == b'\xff\xd8\xff\xe0' or blob[:4] == b'\xff\xd8\xff\xe1' or blob[:4] == b'\xff\xd8\xff\xe2':
                    ext = "jpg"
                elif blob[:8] == b'\x89PNG\r\n\x1a\n':
                    ext = "png"
                else:
                    ext = "jpg"  # default fallback
            filename = f"slide{slide_idx:02d}_img{image_count:03d}.{ext}"
            filepath = os.path.join(IMAGES_DIR, filename)
            with open(filepath, "wb") as f:
                f.write(image.blob)
            shape_info["image_file"] = filename
            shape_info["image_content_type"] = ext
            print(f"  Extracted image: {filename}")

        # Group shapes
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            group_texts = []
            group_images = []
            for grp_shape in shape.shapes:
                if grp_shape.has_text_frame:
                    group_texts.append(grp_shape.text_frame.text)
                if grp_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image = grp_shape.image
                    image_count += 1
                    try:
                        ext = image.content_type.split("/")[-1]
                        if ext == "jpeg":
                            ext = "jpg"
                    except (ValueError, AttributeError):
                        blob = image.blob
                        if blob[:4] == b'\xff\xd8\xff\xe0' or blob[:4] == b'\xff\xd8\xff\xe1' or blob[:4] == b'\xff\xd8\xff\xe2':
                            ext = "jpg"
                        elif blob[:8] == b'\x89PNG\r\n\x1a\n':
                            ext = "png"
                        else:
                            ext = "jpg"
                    filename = f"slide{slide_idx:02d}_img{image_count:03d}.{ext}"
                    filepath = os.path.join(IMAGES_DIR, filename)
                    with open(filepath, "wb") as f:
                        f.write(image.blob)
                    group_images.append(filename)
                    print(f"  Extracted image: {filename}")
            shape_info["group_texts"] = group_texts
            shape_info["group_images"] = group_images

        slide_info["shapes"].append(shape_info)

    slides_data.append(slide_info)
    print(f"Slide {slide_idx}: {slide_info['layout_name']} - {len(slide_info['shapes'])} shapes")

# Save raw extracted data
with open("pptx_raw_data.json", "w") as f:
    json.dump(slides_data, f, indent=2, default=str)

print(f"\nTotal slides: {len(slides_data)}")
print(f"Total images extracted: {image_count}")
print("Raw data saved to pptx_raw_data.json")

# Print readable text summary
print("\n" + "=" * 80)
print("FULL TEXT CONTENT BY SLIDE")
print("=" * 80)

for slide in slides_data:
    print(f"\n--- SLIDE {slide['slide_number']} ({slide['layout_name']}) ---")
    for shape in slide['shapes']:
        if 'text_content' in shape:
            for para in shape['text_content']:
                if para['text'].strip():
                    prefix = "  " * para['level']
                    formatting = []
                    for run in para['runs']:
                        if run['bold']:
                            formatting.append("BOLD")
                        if run['italic']:
                            formatting.append("ITALIC")
                    fmt_str = f" [{', '.join(formatting)}]" if formatting else ""
                    print(f"  {prefix}{para['text']}{fmt_str}")
        if 'table' in shape:
            print(f"  [TABLE]:")
            for row in shape['table']:
                print(f"    | {' | '.join(row)} |")
        if 'image_file' in shape:
            print(f"  [IMAGE: {shape['image_file']}]")
        if 'group_texts' in shape and shape['group_texts']:
            for gt in shape['group_texts']:
                if gt.strip():
                    print(f"  (group) {gt}")
        if 'group_images' in shape and shape['group_images']:
            for gi in shape['group_images']:
                print(f"  [GROUP IMAGE: {gi}]")

    if slide['notes']:
        print(f"  [NOTES: {slide['notes']}]")
